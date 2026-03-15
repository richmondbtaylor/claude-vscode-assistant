"""
Build promptAnything.io Google Slides deck (.pptx → Drive → Google Slides)
Bishop AI brand: Deep Navy bg, Gold + Blue accents, Poppins/Montserrat/Open Sans
Widescreen 16:9 — 13.33" × 7.5"
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
import copy

# ── Colors ────────────────────────────────────────────────────────────────────
DN  = RGBColor(0,   8,   20)   # Deep Navy
DK  = RGBColor(30,  35,  51)   # Dark Navy
OW  = RGBColor(230, 226, 222)  # Off-White
WH  = RGBColor(250, 251, 250)  # White
GO  = RGBColor(224, 184, 72)   # Gold
BL  = RGBColor(24,  148, 201)  # Blue
RED = RGBColor(224, 82,  82)   # Red accent

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=DN):
    """Full-slide background rectangle."""
    sh = slide.shapes.add_shape(1, 0, 0, SW, SH)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(1)):
    """Add a rectangle with optional fill and border."""
    sh = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh


def txb(slide, text, x, y, w, h,
        font="Open Sans", size=Pt(11), color=OW,
        bold=False, align=PP_ALIGN.LEFT, wrap=True,
        italic=False, spacing_after=0):
    """Add a text box."""
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_after:
        p.space_after = Pt(spacing_after)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tf


def label(slide, text, x, y, w=Inches(4), color=BL):
    """Slide section label — small caps Montserrat."""
    txb(slide, text.upper(), x, y, w, Inches(0.3),
        font="Montserrat", size=Pt(8), color=color, bold=True)


def gold_bar(slide, x, y, w=Inches(0.55), h=Pt(3)):
    """Thin gold accent bar."""
    r = rect(slide, x, y, w, h, fill=GO)
    return r


def slide_num(slide, n):
    txb(slide, f"{n:02d}", Inches(0.5), SH - Inches(0.38),
        Inches(0.5), Inches(0.3), font="Montserrat", size=Pt(7), color=OW)


def watermark(slide):
    txb(slide, "promptAnything.io", SW - Inches(2.2), SH - Inches(0.38),
        Inches(2), Inches(0.3), font="Montserrat", size=Pt(7), color=OW,
        align=PP_ALIGN.RIGHT)


def screenshot_placeholder(slide, x, y, w, h, label_text):
    """Dashed-border placeholder for a screenshot."""
    r = rect(slide, x, y, w, h, fill=DK, line_color=GO, line_w=Pt(1.5))
    # inner label
    tf = slide.shapes.add_textbox(x, y + h/2 - Inches(0.4), w, Inches(0.8)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ Screenshot: {label_text} ]"
    run.font.name = "Montserrat"
    run.font.size = Pt(9)
    run.font.color.rgb = GO
    run.font.bold = True
    return r


def card(slide, x, y, w, h, fill=DK, accent=GO, accent_side="left",
         accent_w=Pt(3)):
    """Rounded-corner card with accent edge."""
    r = rect(slide, x, y, w, h, fill=fill)
    if accent_side == "left":
        rect(slide, x, y, accent_w, h, fill=accent)
    elif accent_side == "top":
        rect(slide, x, y, w, accent_w, fill=accent)
    return r


def bullet_list(slide, items, x, y, w, h, color=OW, size=Pt(11), dot_color=GO):
    """Manually drawn bullet list."""
    dy = Inches(0.0)
    line_h = Inches(0.32)
    for item in items:
        # gold dot
        rect(slide, x, y + dy + Inches(0.09), Pt(5), Pt(5), fill=dot_color)
        txb(slide, item, x + Inches(0.22), y + dy, w - Inches(0.22), line_h,
            font="Open Sans", size=size, color=color)
        dy += line_h + Inches(0.05)


def two_col_layout(left_w_frac=0.58):
    """Return (lx, ly, lw, lh, rx, ry, rw, rh) for a standard two-col layout."""
    pad_x = Inches(0.75)
    pad_y = Inches(0.95)
    gap   = Inches(0.4)
    total_w = SW - 2 * pad_x
    lw = total_w * left_w_frac
    rw = total_w * (1 - left_w_frac) - gap
    lx = pad_x
    rx = pad_x + lw + gap
    ly = ry = pad_y
    lh = rh = SH - pad_y - Inches(0.5)
    return lx, ly, lw, lh, rx, ry, rw, rh


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

# ── S1 Title ──────────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
# Subtle accent circle top-right
circ = s.shapes.add_shape(9, SW - Inches(3.2), -Inches(1.5), Inches(5), Inches(5))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(20, 25, 40)
circ.line.fill.background()

label(s, "Webinar · 2026", Inches(0.75), Inches(1.5))
gold_bar(s, Inches(0.75), Inches(1.9), Inches(0.55))
txb(s, "The Universal Interface\nfor Artificial Intelligence.",
    Inches(0.75), Inches(2.05), Inches(8.2), Inches(2.5),
    font="Poppins", size=Pt(48), color=WH, bold=True)
txb(s, "Let's translate your ideas and unleash the true power of AI.",
    Inches(0.75), Inches(4.6), Inches(7), Inches(0.5),
    font="Open Sans", size=Pt(14), color=OW)
txb(s, "promptAnything.io", Inches(0.75), Inches(5.35), Inches(4), Inches(0.4),
    font="Poppins", size=Pt(16), color=GO, bold=True)
txb(s, "The Universal Interface for AI  |  Five Steps. Any Idea. Any Platform.", Inches(2.55), Inches(5.38), Inches(6), Inches(0.35),
    font="Montserrat", size=Pt(8), color=OW)
slide_num(s, 1); watermark(s)

# ── S2 Exec Summary ───────────────────────────────────────────────────────────
s = new_slide(); bg(s)
label(s, "At a Glance", Inches(0.75), Inches(0.55))
txb(s, "95% of AI users are only scratching the surface. The bottleneck isn't the technology. It's us.",
    Inches(0.75), Inches(0.95), Inches(11), Inches(0.65),
    font="Poppins", size=Pt(22), color=WH, bold=True)
gold_bar(s, Inches(0.75), Inches(1.7), Inches(11.5))

# Stat bar
txb(s, "12x LOSS FACTOR.", Inches(0.75), Inches(1.85), Inches(2.1), Inches(0.35),
    font="Montserrat", size=Pt(10), color=GO, bold=True)
txb(s, "Poor prompts create a 12x loss in initial iteration time. promptAnything.io closes that gap in one session.",
    Inches(2.95), Inches(1.85), Inches(9.5), Inches(0.35),
    font="Open Sans", size=Pt(10), color=OW)

# Three cards
CW = Inches(3.8); CX = [Inches(0.75), Inches(4.75), Inches(8.75)]
CY = Inches(2.45); CH = Inches(3.8)
titles  = ["The Problem",    "The Solution",      "The Result"]
accents = [RED,               GO,                  BL]
items   = [
    ["95% of AI users only scratching the surface", "12x loss factor from poor prompt structure", "$270K+ prompt engineers — or you guess alone"],
    ["5-step guided prompt builder", "Platform asks the questions you don't know to ask", "AI-agnostic — works with any model"],
    ["Deploy-ready prompt, one session", "Works in n8n, Make, Claude, and more", "No iterations. No guesswork."],
]
for i in range(3):
    card(s, CX[i], CY, CW, CH, fill=DK, accent=accents[i], accent_side="top", accent_w=Pt(3))
    txb(s, titles[i].upper(), CX[i] + Inches(0.2), CY + Inches(0.18), CW - Inches(0.3), Inches(0.28),
        font="Montserrat", size=Pt(8), color=accents[i], bold=True)
    bullet_list(s, items[i], CX[i] + Inches(0.12), CY + Inches(0.6), CW - Inches(0.25), CH - Inches(0.7),
                dot_color=accents[i], size=Pt(10))
slide_num(s, 2); watermark(s)

# ── S3 Problem ────────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.56)
label(s, "The Problem", lx, Inches(0.55))
gold_bar(s, lx, Inches(0.95))
txb(s, "AI's biggest bottleneck\nisn't the technology. It's us.",
    lx, Inches(1.05), lw, Inches(1.6), font="Poppins", size=Pt(28), color=WH, bold=True)
bullet_list(s, [
    "95% of AI users only scratching the surface",
    "Poor prompts create a 12x loss in iteration time",
    "Static prompts incur a 156% opportunity cost vs. optimized",
    "Prompt engineers earn $270K+ — most can't afford one",
], lx, Inches(2.65), lw, Inches(1.4), size=Pt(11))
txb(s, "It's not the AI's fault. It's our inability to communicate with it. The gap between your idea and a working AI product is filled with guesswork — and that gap has a real cost.",
    lx, Inches(4.15), lw, Inches(1.2), font="Open Sans", size=Pt(10), color=OW)

# Right — cost breakdown
card(s, rx, ry + Inches(0.1), rw, Inches(1.1), fill=DK, accent=RED)
txb(s, "12x", rx + Inches(0.2), ry + Inches(0.15), Inches(0.85), Inches(0.75),
    font="Poppins", size=Pt(28), color=RED, bold=True)
txb(s, "loss factor\nin iteration time", rx + Inches(1.15), ry + Inches(0.22), rw - Inches(1.3), Inches(0.6),
    font="Open Sans", size=Pt(10), color=OW)

card(s, rx, ry + Inches(1.35), rw, Inches(1.1), fill=DK, accent=RED)
txb(s, "156%", rx + Inches(0.2), ry + Inches(1.4), Inches(1.1), Inches(0.75),
    font="Poppins", size=Pt(24), color=RED, bold=True)
txb(s, "opportunity cost\nfrom static prompts", rx + Inches(1.45), ry + Inches(1.47), rw - Inches(1.6), Inches(0.6),
    font="Open Sans", size=Pt(10), color=OW)

card(s, rx, ry + Inches(2.6), rw, Inches(1.1), fill=DK, accent=GO)
txb(s, "$270K+", rx + Inches(0.2), ry + Inches(2.65), Inches(1.3), Inches(0.75),
    font="Poppins", size=Pt(22), color=GO, bold=True)
txb(s, "avg. prompt engineer\nsalary — if you can hire one", rx + Inches(1.65), ry + Inches(2.72), rw - Inches(1.8), Inches(0.65),
    font="Open Sans", size=Pt(10), color=OW)

card(s, rx, ry + Inches(3.85), rw, Inches(1.1), fill=DK, accent=BL)
txb(s, "95%", rx + Inches(0.2), ry + Inches(3.9), Inches(0.95), Inches(0.75),
    font="Poppins", size=Pt(28), color=BL, bold=True)
txb(s, "of AI users are scratching\nthe surface of what's possible", rx + Inches(1.25), ry + Inches(3.97), rw - Inches(1.4), Inches(0.65),
    font="Open Sans", size=Pt(10), color=OW)
slide_num(s, 3); watermark(s)

# ── S4 Solution ───────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.52)
label(s, "The Solution", lx, Inches(0.55))
gold_bar(s, lx, Inches(0.95))
txb(s, "The AI conversation partner\nthat understands your intent.",
    lx, Inches(1.05), lw, Inches(1.8), font="Poppins", size=Pt(26), color=WH, bold=True)
bullet_list(s, [
    "AI-agnostic — works with GPT-4, Claude, Gemini, and any model",
    "Translates natural language goals into precision-engineered prompts",
    "Not a template. An intelligent, guided conversation that builds your prompt",
], lx, Inches(2.85), lw, Inches(1.1), size=Pt(11))
txb(s, "We are not another prompt-fixing tool. We are the bridge over the gap between your idea and the AI's perfect execution — your personal AI strategist.",
    lx, Inches(4.1), lw, Inches(1.4), font="Open Sans", size=Pt(10), color=OW)

# Right — Before/After cards
card(s, rx, ry + Inches(0.1), rw, Inches(1.75), fill=RGBColor(40, 20, 20), accent=RED, accent_side="left")
txb(s, "WITHOUT PROMPTANYTHING", rx + Inches(0.2), ry + Inches(0.25), rw - Inches(0.3), Inches(0.25),
    font="Montserrat", size=Pt(8), color=RED, bold=True)
txb(s, "Blank page. 12x iteration loss. Paying $270K+ for a prompt engineer or giving up entirely. 95% never unlock the model's real potential.",
    rx + Inches(0.2), ry + Inches(0.6), rw - Inches(0.3), Inches(1.05),
    font="Open Sans", size=Pt(10), color=OW)

card(s, rx, ry + Inches(2.1), rw, Inches(1.75), fill=RGBColor(20, 30, 15), accent=GO, accent_side="left")
txb(s, "WITH PROMPTANYTHING", rx + Inches(0.2), ry + Inches(2.25), rw - Inches(0.3), Inches(0.25),
    font="Montserrat", size=Pt(8), color=GO, bold=True)
txb(s, "5 guided steps. One session. A structured, deploy-ready prompt — AI-agnostic, portable, and built around your specific goal.",
    rx + Inches(0.2), ry + Inches(2.6), rw - Inches(0.3), Inches(1.05),
    font="Open Sans", size=Pt(10), color=OW)
slide_num(s, 4); watermark(s)

# ── S5 Workflow Overview ──────────────────────────────────────────────────────
s = new_slide(); bg(s)
label(s, "The promptAnything Workflow", Inches(0.75), Inches(0.55))
txb(s, "Five steps. Every idea. Every platform.",
    Inches(0.75), Inches(0.95), Inches(11), Inches(0.6),
    font="Poppins", size=Pt(28), color=WH, bold=True)
gold_bar(s, Inches(0.75), Inches(1.65), Inches(11.5))

steps = [
    ("01", "IDEA",     "Start with a simple idea — one sentence is enough", GO),
    ("02", "PATH",     "Choose agent workflow or Claude skill", GO),
    ("03", "REFINE",   "Platform asks and answers the right questions", GO),
    ("04", "GENERATE", "One click builds the full structured framework", GO),
    ("05", "DEPLOY",   "Paste into any tool. Fully portable output.", BL),
]
CW = Inches(2.3); GAP = Inches(0.18)
start_x = Inches(0.75)
for i, (num, name, desc, col) in enumerate(steps):
    cx = start_x + i * (CW + GAP)
    card(s, cx, Inches(1.95), CW, Inches(4.85), fill=DK, accent=col, accent_side="top", accent_w=Pt(3))
    txb(s, num, cx + Inches(0.2), Inches(2.15), CW - Inches(0.3), Inches(0.8),
        font="Poppins", size=Pt(36), color=col, bold=True)
    txb(s, name, cx + Inches(0.2), Inches(2.95), CW - Inches(0.3), Inches(0.3),
        font="Montserrat", size=Pt(9), color=WH, bold=True)
    txb(s, desc, cx + Inches(0.2), Inches(3.35), CW - Inches(3), Inches(0.9),
        font="Open Sans", size=Pt(10), color=OW)
    # arrow between cards
    if i < 4:
        ax = cx + CW + Inches(0.04)
        ay = Inches(4.25)
        txb(s, "→", ax, ay, GAP + Inches(0.1), Inches(0.35),
            font="Poppins", size=Pt(14), color=GO, align=PP_ALIGN.CENTER)
slide_num(s, 5); watermark(s)

# ── S6 Step 1: IDEA ───────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.5)
label(s, "Step 01 of 05  ·  IDEA", lx, Inches(0.55), color=GO)
gold_bar(s, lx, Inches(0.95))
txb(s, "You already have everything\nstep one needs.",
    lx, Inches(1.05), lw, Inches(1.7), font="Poppins", size=Pt(28), color=WH, bold=True)
bullet_list(s, [
    "One sentence minimum",
    "No technical framing required",
    "Any domain: business, content, automation",
], lx, Inches(2.75), lw, Inches(1.1), size=Pt(11))
card(s, lx, Inches(3.95), lw, Inches(1.5), fill=DK, accent=GO)
txb(s, '"I want an AI that screens my inbound leads before I get on a call with them."',
    lx + Inches(0.22), Inches(4.1), lw - Inches(0.35), Inches(0.85),
    font="Open Sans", size=Pt(11), color=OW, italic=True)
txb(s, "That's all step one needs.",
    lx + Inches(0.22), Inches(5.05), lw - Inches(0.35), Inches(0.3),
    font="Montserrat", size=Pt(9), color=GO, bold=True)

screenshot_placeholder(s, rx, ry, rw, rh, "Idea input screen")
slide_num(s, 6); watermark(s)

# ── S7 Step 2: PATH ───────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.5)
label(s, "Step 02 of 05  ·  PATH", lx, Inches(0.55), color=GO)
gold_bar(s, lx, Inches(0.95))
txb(s, "Two paths. One decision.\nNo wrong answer.",
    lx, Inches(1.05), lw, Inches(1.7), font="Poppins", size=Pt(28), color=WH, bold=True)
bullet_list(s, [
    "Agent workflow: multi-step, tool-connected",
    "Claude skill: single-purpose, conversation-embedded",
    "The choice shapes the entire prompt structure",
], lx, Inches(2.75), lw, Inches(1.1), size=Pt(11))
txb(s, "Most people never think about this distinction before writing. That single decision changes everything about the prompt that gets built.",
    lx, Inches(3.95), lw, Inches(1.0), font="Open Sans", size=Pt(10), color=OW)

# Path cards on left
card(s, lx, Inches(5.1), lw * 0.47, Inches(1.4), fill=DK, accent=GO)
txb(s, "AGENT WORKFLOW", lx + Inches(0.2), Inches(5.25), lw * 0.47 - Inches(0.25), Inches(0.25),
    font="Montserrat", size=Pt(8), color=GO, bold=True)
txb(s, "n8n · Make · Zapier\nBest for: multi-step automation, tool connections",
    lx + Inches(0.2), Inches(5.55), lw * 0.47 - Inches(0.25), Inches(0.8),
    font="Open Sans", size=Pt(9), color=OW)

card(s, lx + lw * 0.47 + Inches(0.18), Inches(5.1), lw * 0.47, Inches(1.4), fill=DK, accent=BL)
txb(s, "CLAUDE SKILL", lx + lw * 0.47 + Inches(0.38), Inches(5.25), lw * 0.47 - Inches(0.25), Inches(0.25),
    font="Montserrat", size=Pt(8), color=BL, bold=True)
txb(s, "Repeatable · Single-purpose\nBest for: consistent reusable behavior in Claude",
    lx + lw * 0.47 + Inches(0.38), Inches(5.55), lw * 0.47 - Inches(0.25), Inches(0.8),
    font="Open Sans", size=Pt(9), color=OW)

screenshot_placeholder(s, rx, ry, rw, rh, "Path selection screen")
slide_num(s, 7); watermark(s)

# ── S8 Step 3: REFINE ─────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.5)
label(s, "Step 03 of 05  ·  REFINE", lx, Inches(0.55), color=GO)
gold_bar(s, lx, Inches(0.95))
txb(s, "The platform asks what you\nforgot to think about.",
    lx, Inches(1.05), lw, Inches(1.7), font="Poppins", size=Pt(28), color=WH, bold=True)
bullet_list(s, [
    "Relevant questions auto-generated to your idea",
    "Fills in context, constraints, and logic gaps",
    "No prompt engineering knowledge required",
], lx, Inches(2.75), lw, Inches(1.1), size=Pt(11))
txb(s, "It doesn't just ask — it answers alongside you. Compressed into a guided conversation anyone can complete.",
    lx, Inches(3.95), lw, Inches(0.75), font="Open Sans", size=Pt(10), color=OW)

# Sample questions
qs = [
    "Who is the audience for this output?",
    "What format should the AI use?",
    "What should it never do?",
    "What's the success condition?",
]
for i, q in enumerate(qs):
    qy = Inches(4.85) + i * Inches(0.5)
    card(s, lx, qy, lw, Inches(0.38), fill=DK, accent=GO)
    txb(s, q, lx + Inches(0.2), qy + Inches(0.07), lw - Inches(0.3), Inches(0.28),
        font="Open Sans", size=Pt(10), color=OW)

screenshot_placeholder(s, rx, ry, rw, rh, "Refine Q&A interface")
slide_num(s, 8); watermark(s)

# ── S9 Step 4: GENERATE ───────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.5)
label(s, "Step 04 of 05  ·  GENERATE", lx, Inches(0.55), color=GO)
gold_bar(s, lx, Inches(0.95))
txb(s, "One click builds the\nfull structured framework.",
    lx, Inches(1.05), lw, Inches(1.7), font="Poppins", size=Pt(28), color=WH, bold=True)
bullet_list(s, [
    "Complete prompt generated from your refined idea",
    "Structured format — not freeform text",
    "Deploy-ready on output, no editing required",
], lx, Inches(2.75), lw, Inches(1.1), size=Pt(11))
txb(s, "A framework with context blocks, instruction sets, constraints, and output format specs. Everything a prompt needs to work consistently.",
    lx, Inches(3.95), lw, Inches(0.9), font="Open Sans", size=Pt(10), color=OW)

# Generate button + output preview
btn = rect(s, lx, Inches(5.05), Inches(2), Inches(0.55), fill=GO)
txb(s, "Generate Prompt", lx + Inches(0.15), Inches(5.13), Inches(1.7), Inches(0.38),
    font="Poppins", size=Pt(11), color=DN, bold=True)
txb(s, "↓", lx + Inches(0.8), Inches(5.68), Inches(0.5), Inches(0.3),
    font="Poppins", size=Pt(14), color=GO, align=PP_ALIGN.CENTER)
card(s, lx, Inches(6.05), lw, Inches(1.05), fill=DK, accent=GO, accent_side="left")
txb(s, "[CONTEXT]  You are a...\n[INSTRUCTION]  When given...\n[CONSTRAINTS]  Never...\n[OUTPUT FORMAT]  Return as...",
    lx + Inches(0.22), Inches(6.12), lw - Inches(0.35), Inches(0.9),
    font="Open Sans", size=Pt(9), color=OW)

screenshot_placeholder(s, rx, ry, rw, rh, "Generated prompt output")
slide_num(s, 9); watermark(s)

# ── S10 Step 5: DEPLOY ────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.5)
label(s, "Step 05 of 05  ·  DEPLOY", lx, Inches(0.55), color=BL)
rect(s, lx, Inches(0.95), Inches(0.55), Pt(3), fill=BL)
txb(s, "The prompt works in any\ntool you already use.",
    lx, Inches(1.05), lw, Inches(1.7), font="Poppins", size=Pt(28), color=WH, bold=True)
bullet_list(s, [
    "Drop into n8n, Make, Zapier, or any agent",
    "Use in Claude skills, web apps, custom workflows",
    "One prompt. Total portability.",
], lx, Inches(2.75), lw, Inches(1.1), size=Pt(11))
txb(s, "You're not locked into our platform. You're building assets you own. Every new platform becomes another deployment target.",
    lx, Inches(3.95), lw, Inches(0.9), font="Open Sans", size=Pt(10), color=OW)

# Platform grid
platforms = [("n8n", GO), ("Make", GO), ("Zapier", GO),
             ("Claude Skills", BL), ("ChatGPT", BL), ("Custom Agents", BL)]
GW = lw / 3 - Inches(0.08)
for i, (name, col) in enumerate(platforms):
    px = lx + (i % 3) * (GW + Inches(0.12))
    py = Inches(4.98) + (i // 3) * Inches(0.65)
    card(s, px, py, GW, Inches(0.52), fill=DK, accent=col, accent_side="top", accent_w=Pt(2))
    txb(s, name, px + Inches(0.12), py + Inches(0.12), GW - Inches(0.18), Inches(0.3),
        font="Montserrat", size=Pt(9), color=col, bold=True)

screenshot_placeholder(s, rx, ry, rw, rh, "Deploy — tool destinations")
slide_num(s, 10); watermark(s)

# ── S11 Who It's For ──────────────────────────────────────────────────────────
s = new_slide(); bg(s)
label(s, "Who It's For", Inches(0.75), Inches(0.55))
txb(s, "Built for builders who are not yet technical.",
    Inches(0.75), Inches(0.95), Inches(11), Inches(0.7),
    font="Poppins", size=Pt(28), color=WH, bold=True)
gold_bar(s, Inches(0.75), Inches(1.75), Inches(11.5))

profiles = [
    ("Founders", GO,
     "AI is on the roadmap. Building it keeps hitting a wall at the prompt layer. promptAnything.io removes that wall."),
    ("Operators", GO,
     "Five automation ideas, three months of Zapier experience. Knows the tools but not how to structure a prompt that controls the output."),
    ("Agencies", GO,
     "Delivering AI solutions to clients. Cannot spend eight hours per project on prompts through trial and error."),
]
CW = Inches(3.9); CY = Inches(2.05); CH = Inches(4.55)
for i, (title, col, desc) in enumerate(profiles):
    cx = Inches(0.75) + i * (CW + Inches(0.25))
    card(s, cx, CY, CW, CH, fill=DK, accent=col, accent_side="top", accent_w=Pt(3))
    txb(s, title.upper(), cx + Inches(0.2), CY + Inches(0.25), CW - Inches(0.3), Inches(0.28),
        font="Montserrat", size=Pt(9), color=col, bold=True)
    txb(s, desc, cx + Inches(0.2), CY + Inches(0.7), CW - Inches(0.3), Inches(3.5),
        font="Open Sans", size=Pt(11), color=OW)
slide_num(s, 11); watermark(s)

# ── S12 Market ────────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.55)
label(s, "Market Opportunity", lx, Inches(0.55))
gold_bar(s, lx, Inches(0.95))
txb(s, "Prompt engineering is the\nnew interface layer.",
    lx, Inches(1.05), lw, Inches(1.7), font="Poppins", size=Pt(28), color=WH, bold=True)
bullet_list(s, [
    "Billions in AI spend failing at the prompt layer",
    "Non-technical builders are the fastest-growing AI user segment",
    "Platform-agnostic tools compound across every AI deployment",
], lx, Inches(2.75), lw, Inches(1.1), size=Pt(11))
txb(s, "The model can do anything — the prompt is the constraint. Every new platform is a new deployment target. The market grows as AI grows.",
    lx, Inches(3.95), lw, Inches(0.9), font="Open Sans", size=Pt(10), color=OW)

# Three stat cards — real data from seed deck
stats = [
    ("$32T",  GO,  "Global AI Market by 2034 (TAM)", "The entire ecosystem expands — every new platform is a new deployment target"),
    ("1B+",   GO,  "Knowledge workers (SAM)", "Directly benefit from better AI communication tools — multi-hundred-billion annual market"),
    ("10M",   BL,  "Initial SOM target", "Tech-forward professionals, marketers, developers, and creators actively adopting AI tools"),
]
for i, (num, col, title, sub) in enumerate(stats):
    sy = ry + i * (Inches(1.65))
    card(s, rx, sy, rw, Inches(1.5), fill=DK, accent=col)
    txb(s, num, rx + Inches(0.22), sy + Inches(0.15), Inches(1.2), Inches(0.85),
        font="Poppins", size=Pt(30), color=col, bold=True)
    txb(s, title.upper(), rx + Inches(1.55), sy + Inches(0.18), rw - Inches(1.7), Inches(0.25),
        font="Montserrat", size=Pt(8), color=col, bold=True)
    txb(s, sub, rx + Inches(1.55), sy + Inches(0.52), rw - Inches(1.7), Inches(0.85),
        font="Open Sans", size=Pt(9.5), color=OW)
slide_num(s, 12); watermark(s)

# ── S13 Why Now ───────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.55)
label(s, "Why Now", lx, Inches(0.55))
gold_bar(s, lx, Inches(0.95))
txb(s, "The market is seeking a universal\ntranslator for AI. That is promptAnything.io.",
    lx, Inches(1.05), lw, Inches(1.7), font="Poppins", size=Pt(24), color=WH, bold=True)
bullet_list(s, [
    "Cambrian explosion of AI — $32T market by 2034",
    "Hundreds of specialized models, zero universal interface",
    "Choice overload is stalling enterprise AI adoption",
], lx, Inches(2.75), lw, Inches(1.1), size=Pt(11))

reasons = [
    ("1", "Cambrian Explosion of AI. The market touches nearly every aspect of the global economy. The wave is here."),
    ("2", "Fragmentation. Hundreds of specialized models emerging — nobody knows which to use or how to direct any of them."),
    ("3", "No universal translator existed. Until promptAnything.io."),
]
for i, (n, text) in enumerate(reasons):
    ry2 = ry + i * Inches(1.55)
    col = GO if n == "3" else DK
    card(s, rx, ry2, rw, Inches(1.35), fill=col if n == "3" else DK,
         accent=GO if n == "3" else None)
    txb(s, n, rx + Inches(0.18), ry2 + Inches(0.15), Inches(0.4), Inches(0.65),
        font="Poppins", size=Pt(22), color=GO, bold=True)
    txb(s, text, rx + Inches(0.7), ry2 + Inches(0.25), rw - Inches(0.85), Inches(0.85),
        font="Open Sans", size=Pt(10),
        color=GO if n == "3" else OW,
        bold=(n == "3"))
slide_num(s, 13); watermark(s)

# ── S14 Use Case 1 ────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
label(s, "Use Case · Agent Workflow Automation", Inches(0.75), Inches(0.55), color=GO)
txb(s, "From idea to deployed lead qualifier in one session.",
    Inches(0.75), Inches(0.95), Inches(11.5), Inches(0.7),
    font="Poppins", size=Pt(28), color=WH, bold=True)
gold_bar(s, Inches(0.75), Inches(1.75), Inches(11.5))

cols = [
    ("THE IDEA", GO, '"Screen my inbound leads before I get on a call with them."', True),
    ("THE RESULT", GO, "Deployed in n8n same day\nConnected to intake form + CRM\nZero iterations after generate", False),
    ("TIME TO DEPLOY", GO, "15 min\nidea to running workflow", False),
]
CW2 = Inches(3.9); CY2 = Inches(2.1); CH2 = Inches(4.7)
for i, (title, col, body, italic) in enumerate(cols):
    cx = Inches(0.75) + i * (CW2 + Inches(0.22))
    card(s, cx, CY2, CW2, CH2, fill=DK, accent=col, accent_side="top", accent_w=Pt(3))
    txb(s, title, cx + Inches(0.2), CY2 + Inches(0.22), CW2 - Inches(0.3), Inches(0.28),
        font="Montserrat", size=Pt(8), color=col, bold=True)
    if title == "TIME TO DEPLOY":
        txb(s, "15", cx + Inches(0.5), CY2 + Inches(0.7), CW2 - Inches(0.7), Inches(1.4),
            font="Poppins", size=Pt(55), color=GO, bold=True, align=PP_ALIGN.CENTER)
        txb(s, "min\nidea to running workflow", cx + Inches(0.2), CY2 + Inches(2.1), CW2 - Inches(0.3), Inches(0.7),
            font="Open Sans", size=Pt(10), color=OW, align=PP_ALIGN.CENTER)
    else:
        txb(s, body, cx + Inches(0.2), CY2 + Inches(0.7), CW2 - Inches(0.3), Inches(3.6),
            font="Open Sans", size=Pt(11), color=OW, italic=italic)
slide_num(s, 14); watermark(s)

# ── S15 Use Case 2 ────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
label(s, "Use Case · Claude Skill", Inches(0.75), Inches(0.55), color=BL)
txb(s, "A client-facing content skill built in one afternoon.",
    Inches(0.75), Inches(0.95), Inches(11.5), Inches(0.7),
    font="Poppins", size=Pt(28), color=WH, bold=True)
rect(s, Inches(0.75), Inches(1.75), Inches(11.5), Pt(3), fill=BL)

cols2 = [
    ("THE IDEA", BL, '"Turn any meeting transcript into a follow-up email in my client\'s brand voice."', True),
    ("THE RESULT", BL, "Loaded into Claude skill system\nConsistent tone every invocation\nClient-ready output, no editing", False),
    ("TIME TO DEPLOY", BL, "", False),
]
for i, (title, col, body, italic) in enumerate(cols2):
    cx = Inches(0.75) + i * (CW2 + Inches(0.22))
    card(s, cx, CY2, CW2, CH2, fill=DK, accent=col, accent_side="top", accent_w=Pt(3))
    txb(s, title, cx + Inches(0.2), CY2 + Inches(0.22), CW2 - Inches(0.3), Inches(0.28),
        font="Montserrat", size=Pt(8), color=col, bold=True)
    if title == "TIME TO DEPLOY":
        txb(s, "15", cx + Inches(0.5), CY2 + Inches(0.7), CW2 - Inches(0.7), Inches(1.4),
            font="Poppins", size=Pt(55), color=BL, bold=True, align=PP_ALIGN.CENTER)
        txb(s, "min\nidea to repeatable skill", cx + Inches(0.2), CY2 + Inches(2.1), CW2 - Inches(0.3), Inches(0.7),
            font="Open Sans", size=Pt(10), color=OW, align=PP_ALIGN.CENTER)
    else:
        txb(s, body, cx + Inches(0.2), CY2 + Inches(0.7), CW2 - Inches(0.3), Inches(3.6),
            font="Open Sans", size=Pt(11), color=OW, italic=italic)
slide_num(s, 15); watermark(s)

# ── S16 The Ask ───────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
lx, ly, lw, lh, rx, ry, rw, rh = two_col_layout(0.52)
label(s, "The Ask", lx, Inches(0.55))
gold_bar(s, lx, Inches(0.95))
txb(s, "You already have the idea.\nYou know where the gap is.",
    lx, Inches(1.05), lw, Inches(1.7), font="Poppins", size=Pt(28), color=WH, bold=True)
bullet_list(s, [
    "The prompt is the only thing between you and shipping",
    "You now have a five-step system to close that gap",
    "One session. One prompt. Everything changes.",
], lx, Inches(2.75), lw, Inches(1.1), size=Pt(11))
txb(s, "The one AI idea you've been sitting on — half-built in your head for weeks — is one session away from being real.",
    lx, Inches(3.95), lw, Inches(0.9), font="Open Sans", size=Pt(10), color=OW)

# Outcome box
card(s, lx, Inches(4.95), lw, Inches(1.85), fill=RGBColor(18, 22, 12),
     accent=GO, accent_side="top", accent_w=Pt(2))
txb(s, "WHAT YOU WALK AWAY WITH", lx + Inches(0.2), Inches(5.12), lw - Inches(0.3), Inches(0.25),
    font="Montserrat", size=Pt(8), color=GO, bold=True)
bullet_list(s, [
    "A structured, deploy-ready prompt",
    "Zero iterations required after generate",
    "Works in any tool you already use",
    "Repeatable — run again on any new idea",
], lx + Inches(0.05), Inches(5.45), lw - Inches(0.1), Inches(1.2), size=Pt(10))

# CTA card
card(s, rx, ry + Inches(0.5), rw, Inches(1.5), fill=DK, accent=GO, accent_side="top", accent_w=Pt(3))
txb(s, "No account needed to start.", rx + Inches(0.2), ry + Inches(0.75), rw - Inches(0.3), Inches(0.4),
    font="Poppins", size=Pt(14), color=WH, bold=True)
txb(s, "promptAnything.io → open → build",
    rx + Inches(0.2), ry + Inches(1.2), rw - Inches(0.3), Inches(0.35),
    font="Montserrat", size=Pt(10), color=GO, bold=True)

# Step summary
step_labels = [("1","IDEA","go"),("2","PATH","go"),("3","REFINE","go"),("4","GENERATE","go"),("5","DEPLOY","bl")]
for i, (n, lbl, col) in enumerate(step_labels):
    sx = rx + Inches(0.1) + i * Inches(1.1)
    sy2 = ry + Inches(2.7)
    dot_col = GO if col == "go" else BL
    circ2 = s.shapes.add_shape(9, sx, sy2, Inches(0.36), Inches(0.36))
    circ2.fill.solid(); circ2.fill.fore_color.rgb = dot_col
    circ2.line.fill.background()
    txb(s, n, sx + Inches(0.07), sy2 + Inches(0.04), Inches(0.24), Inches(0.28),
        font="Poppins", size=Pt(11), color=DN, bold=True, align=PP_ALIGN.CENTER)
    txb(s, lbl, sx - Inches(0.05), sy2 + Inches(0.45), Inches(0.48), Inches(0.3),
        font="Montserrat", size=Pt(7), color=dot_col, align=PP_ALIGN.CENTER)
    if i < 4:
        txb(s, "→", sx + Inches(0.45), sy2 + Inches(0.06), Inches(0.35), Inches(0.28),
            font="Poppins", size=Pt(11), color=GO)
slide_num(s, 16); watermark(s)

# ── S17 Closing ───────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
# large decorative circle
circ3 = s.shapes.add_shape(9, SW - Inches(4), -Inches(2), Inches(7), Inches(7))
circ3.fill.solid(); circ3.fill.fore_color.rgb = RGBColor(15, 20, 35)
circ3.line.fill.background()

label(s, "Start Today", Inches(3), Inches(1.2), color=BL)
txb(s, "Try the workflow on your idea today.",
    Inches(1.5), Inches(1.65), Inches(10), Inches(1.0),
    font="Poppins", size=Pt(36), color=WH, bold=True, align=PP_ALIGN.CENTER)
gold_bar(s, SW/2 - Inches(0.28), Inches(2.75), Inches(0.55))

# CTA box
card(s, Inches(3.2), Inches(3.0), Inches(6.9), Inches(2.4),
     fill=RGBColor(15, 20, 35), accent=GO, accent_side="top", accent_w=Pt(2))
txb(s, "1. Go to promptAnything.io", Inches(3.5), Inches(3.2), Inches(6.3), Inches(0.35),
    font="Montserrat", size=Pt(12), color=GO, bold=True)
txb(s, "Choose your path. Run the workflow.",
    Inches(3.5), Inches(3.6), Inches(6.3), Inches(0.3),
    font="Open Sans", size=Pt(11), color=OW)
txb(s, "2. Walk out with a deployment-ready prompt.",
    Inches(3.5), Inches(3.98), Inches(6.3), Inches(0.35),
    font="Montserrat", size=Pt(12), color=GO, bold=True)
txb(s, "One session. No guesswork. No iterations.",
    Inches(3.5), Inches(4.38), Inches(6.3), Inches(0.3),
    font="Open Sans", size=Pt(11), color=OW)

txb(s, "promptAnything.io",
    Inches(4.5), Inches(5.6), Inches(4.3), Inches(0.55),
    font="Poppins", size=Pt(22), color=GO, bold=True, align=PP_ALIGN.CENTER)

# Workflow bar at bottom
wf_labels = [("IDEA", GO), ("PATH", GO), ("REFINE", GO), ("GENERATE", GO), ("DEPLOY", BL)]
bar_x = Inches(1.2); bar_y = Inches(6.7); seg_w = Inches(2.1)
for i, (lbl, col) in enumerate(wf_labels):
    bx = bar_x + i * (seg_w + Inches(0.08))
    txb(s, lbl, bx, bar_y, seg_w, Inches(0.35),
        font="Montserrat", size=Pt(8), color=col, bold=True, align=PP_ALIGN.CENTER)
    if i < 4:
        rect(s, bx + seg_w, bar_y + Inches(0.1), Inches(0.08), Inches(0.12), fill=GO)

slide_num(s, 17); watermark(s)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "C:/Users/richm/.claude/2026-03-13 promptAnything.io - Keynote Pitch Deck.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
