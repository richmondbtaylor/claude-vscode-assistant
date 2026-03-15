"""
patch_deck.py  v3 — SaaS visual refresh
Modern SaaS aesthetic: rounded cards, pill badges, layered depth,
aggressive blue/gold hierarchy, code-editor demo blocks, progress dots.
"""

import subprocess
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand ─────────────────────────────────────────────────────────────────────
DEEP_NAVY  = RGBColor(0,   8,  20)
DARK_NAVY  = RGBColor(30,  35, 51)
MID_NAVY   = RGBColor(20,  26, 46)      # slightly lighter card surface
CARD       = RGBColor(18,  22, 38)      # card background
GOLD       = RGBColor(224, 184, 72)
BLUE       = RGBColor(24,  148, 201)
BLUE_DIM   = RGBColor(16,  100, 140)    # muted blue for borders
WHITE      = RGBColor(250, 251, 250)
OFF_WHITE  = RGBColor(200, 205, 215)    # slightly cooler for SaaS feel
DIM        = RGBColor(120, 130, 150)    # dim text

F_HEAD = "Poppins"
F_SUB  = "Montserrat"
F_BODY = "Open Sans"
F_MONO = "Courier New"

W = Inches(13.333)
H = Inches(7.5)

SLIDE_TITLES = [
    "", "The Problem", "About Richmond", "The Data",
    "4 Message Types", "5-Step Loop", "PromptAnything",
    "Build 1 Setup", "Build 1 Output", "Build 2 Setup",
    "Build 2 Output", "Build 3 Setup", "Build 3 Output",
    "What You Built", "Book a Call",
]

OUTPUT = Path("c:/Users/richm/.claude") / f"{date.today()} Prompt Engineering for Sales Professionals - Webinar.pptx"
UPLOAD = Path("c:/Users/richm/.claude/skills/infographic-generator/scripts/upload_gdrive.py")


# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs, bg=DEEP_NAVY):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = bg
    return s

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def rgb_hex(c: RGBColor):
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


# ── Shape builders ─────────────────────────────────────────────────────────────

RECT     = 1   # MSO_AUTO_SHAPE_TYPE.RECTANGLE
ROUNDED  = 5   # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE

def shape(slide, kind, l, t, w, h, fill, line_color=None, line_w=0.75, adj=None):
    s = slide.shapes.add_shape(kind, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    # adjust corner radius for rounded rect
    if kind == ROUNDED and adj is not None:
        sp = s._element
        prstGeom = sp.find(qn("a:prstGeom"))
        if prstGeom is not None:
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn("a:avLst"))
            else:
                for gd in avLst.findall(qn("a:gd")):
                    avLst.remove(gd)
            gd = etree.SubElement(avLst, qn("a:gd"))
            gd.set("name", "adj")
            gd.set("fmla", f"val {adj}")
    return s

def rr(slide, l, t, w, h, fill, line_color=None, line_w=0.75, corner=8000):
    """Rounded rectangle with consistent corner radius."""
    return shape(slide, ROUNDED, l, t, w, h, fill, line_color, line_w, adj=corner)

def box(slide, l, t, w, h, text, fn=F_BODY, fs=15, color=WHITE,
        bold=False, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.name      = fn
    r.font.size      = Pt(fs)
    r.font.color.rgb = color
    r.font.bold      = bold
    r.font.italic    = italic
    return tb

def pill(slide, l, t, text, bg=BLUE, fg=WHITE, fs=10):
    """Small rounded pill badge — for labels and tags."""
    w = Inches(len(text) * 0.095 + 0.4)
    h = Inches(0.32)
    rr(slide, l, t, w, h, bg, corner=20000)
    box(slide, l, t + Inches(0.02), w, h - Inches(0.04),
        text, fn=F_SUB, fs=fs, color=fg, bold=True, align=PP_ALIGN.CENTER)
    return w


# ── Chrome ─────────────────────────────────────────────────────────────────────

def logo_wordmark(slide, x=Inches(0.45), y=Inches(0.18)):
    box(slide, x, y, Inches(1.1), Inches(0.38),
        "BISHOP", fn=F_HEAD, fs=15, color=WHITE, bold=True)
    box(slide, x + Inches(1.05), y, Inches(0.5), Inches(0.38),
        "AI", fn=F_HEAD, fs=15, color=GOLD, bold=True)

def footer_bar(slide):
    shape(slide, RECT, Inches(0), Inches(7.22), W, Inches(0.28), CARD)
    box(slide, Inches(0.5), Inches(7.24), Inches(12.3), Inches(0.22),
        "PromptAnything.io  ·  BishopAI.com  ·  calendly.com/richmond  ·  Automation Nation",
        fn=F_BODY, fs=9, color=DIM, align=PP_ALIGN.CENTER)

def progress_dots(slide, current, total=14):
    """Small progress indicator dots, top center."""
    dot_w = Inches(0.12)
    gap   = Inches(0.06)
    total_w = total * dot_w + (total - 1) * gap
    start_x = (W - total_w) / 2
    for i in range(total):
        c = GOLD if i == current - 1 else DIM
        rr(slide,
           start_x + i * (dot_w + gap),
           Inches(0.1), dot_w, Inches(0.12),
           c, corner=20000)

def section_label(slide, text, color=BLUE):
    pill(slide, Inches(0.45), Inches(0.75), text, bg=color, fs=10)

def headline(slide, text, fs=30, color=WHITE, t=Inches(1.18)):
    box(slide, Inches(0.45), t, Inches(12.4), Inches(0.85),
        text, fn=F_HEAD, fs=fs, color=color, bold=True)

def divider(slide, t=Inches(2.1), color=CARD):
    shape(slide, RECT, Inches(0.45), t, Inches(12.4), Inches(0.02), color)

def chrome(slide, n, label=None, label_color=BLUE):
    logo_wordmark(slide)
    footer_bar(slide)
    progress_dots(slide, n)
    if label:
        section_label(slide, label, label_color)


# ── Stat card component ────────────────────────────────────────────────────────

def stat_card(slide, l, t, w, h, stat, label, source=None,
              stat_color=GOLD, border=BLUE_DIM):
    rr(slide, l, t, w, h, CARD, line_color=border, line_w=0.5, corner=6000)
    # Accent top strip
    rr(slide, l, t, w, Inches(0.08), stat_color, corner=6000)
    box(slide, l, t + Inches(0.18), w, Inches(1.1),
        stat, fn=F_HEAD, fs=52, color=stat_color, bold=True, align=PP_ALIGN.CENTER)
    box(slide, l + Inches(0.15), t + Inches(1.35), w - Inches(0.3), Inches(0.65),
        label, fn=F_BODY, fs=13, color=OFF_WHITE, align=PP_ALIGN.CENTER)
    if source:
        pill(slide, l + w/2 - Inches(0.8), t + h - Inches(0.42),
             source, bg=MID_NAVY, fg=DIM, fs=9)


# ── SLIDES ────────────────────────────────────────────────────────────────────

def s_cover(prs):
    slide = blank(prs, DEEP_NAVY)

    # Full-bleed gradient feel: layered rectangles
    shape(slide, RECT, Inches(0), Inches(0), W, H, DEEP_NAVY)
    # Blue glow top-right
    shape(slide, RECT, Inches(8), Inches(0), Inches(5.3), Inches(3.5), MID_NAVY)
    # Gold accent bar left
    shape(slide, RECT, Inches(0), Inches(0), Inches(0.22), H, GOLD)
    # Dark card center
    rr(slide, Inches(1.8), Inches(1.1), Inches(9.7), Inches(5.1), CARD,
       line_color=BLUE_DIM, line_w=0.75, corner=4000)

    # BISHOP AI wordmark inside card
    box(slide, Inches(2.5), Inches(1.35), Inches(2.2), Inches(0.52),
        "BISHOP", fn=F_HEAD, fs=22, color=WHITE, bold=True)
    box(slide, Inches(4.65), Inches(1.35), Inches(0.7), Inches(0.52),
        "AI", fn=F_HEAD, fs=22, color=GOLD, bold=True)
    # Blue pill label
    pill(slide, Inches(5.5), Inches(1.45), "LIVE WEBINAR", bg=BLUE, fs=10)

    # Main title
    box(slide, Inches(2.1), Inches(2.05), Inches(9.3), Inches(1.7),
        "Prompt Engineering\nfor Sales Professionals",
        fn=F_HEAD, fs=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Gold rule
    shape(slide, RECT, Inches(4.2), Inches(3.85), Inches(4.9), Inches(0.04), GOLD)

    # Subtitle
    box(slide, Inches(2.1), Inches(4.0), Inches(9.3), Inches(0.52),
        "The #1 Skill in AI — Learned in 60 Minutes",
        fn=F_SUB, fs=20, color=GOLD, align=PP_ALIGN.CENTER)

    # Presenter
    box(slide, Inches(2.1), Inches(4.65), Inches(9.3), Inches(0.42),
        "Richmond Taylor  ·  Founder, Bishop AI",
        fn=F_BODY, fs=16, color=OFF_WHITE, align=PP_ALIGN.CENTER)

    # Three feature pills at bottom of card
    labels = ["3 Live Builds", "Starts in Claude", "Keep Everything"]
    pill_x = Inches(2.8)
    for lbl in labels:
        w = pill(slide, pill_x, Inches(5.6), lbl, bg=DARK_NAVY, fg=OFF_WHITE, fs=11)
        pill_x += w + Inches(0.25)

    footer_bar(slide)
    notes(slide, "Hold on screen as attendees join. No talking. Start with Hook 2 the moment the session begins.")
    return slide


def s_hook(prs, n):
    slide = blank(prs)
    chrome(slide, n, "THE PROBLEM", GOLD)
    headline(slide, "More reps missed quota last year than the year before.", fs=28)

    stats = [
        ("28%",  "of the week is\nactual selling",          "Salesforce 2025", GOLD),
        ("78%",  "of sellers\nmissed quota",                "Salesmotion 2025", RGBColor(220, 80, 80)),
        ("3.7x", "more likely to hit quota\nwith structured AI",  "Gartner 2025", BLUE),
    ]
    cw = Inches(3.9)
    for i, (stat, label, src, color) in enumerate(stats):
        stat_card(slide, Inches(0.45) + i * (cw + Inches(0.2)),
                  Inches(2.05), cw, Inches(3.5),
                  stat, label, src, stat_color=color)

    box(slide, Inches(0.45), Inches(5.78), Inches(12.4), Inches(0.42),
        "The tools aren't getting worse. Reps aren't getting lazier. The prompts are wrong.",
        fn=F_SUB, fs=14, color=DIM, align=PP_ALIGN.CENTER, italic=True)

    notes(slide, "Hook 2: '78% of sellers missed quota last year. The year before, 69%. So what is actually happening?'\nPause after each stat card. Let them land.")
    return slide


def s_origin(prs, n):
    slide = blank(prs)
    chrome(slide, n, "ABOUT RICHMOND")
    headline(slide, "Pro athlete. Sales director. AI builder.", fs=28)

    items = [
        (GOLD, "⚽  Professional soccer",    "Costa Rica & Spain — then the dream ended, without a plan."),
        (BLUE, "📈  Director of Sales",       "Salesforce ecosystem, consulting, GTM strategy — learned the whole system."),
        (GOLD, "🤖  Built Bishop AI",         "Got obsessed with the gap between prompts that get mediocre output and prompts that get something actually useful."),
    ]
    for i, (accent, title, desc) in enumerate(items):
        t = Inches(2.1) + i * Inches(1.6)
        rr(slide, Inches(0.45), t, Inches(12.4), Inches(1.35), CARD,
           line_color=accent, line_w=0.75, corner=4000)
        shape(slide, RECT, Inches(0.45), t, Inches(0.15), Inches(1.35), accent)
        box(slide, Inches(0.75), t + Inches(0.15), Inches(12.0), Inches(0.45),
            title, fn=F_SUB, fs=15, color=WHITE, bold=True)
        box(slide, Inches(0.75), t + Inches(0.62), Inches(11.5), Inches(0.6),
            desc, fn=F_BODY, fs=13, color=OFF_WHITE)

    pill(slide, Inches(0.45), Inches(7.0), "CHAT: 1–5 — how much AI are you using in your daily sales work?",
         bg=DARK_NAVY, fg=BLUE, fs=11)

    notes(slide, "Tell the full story honestly. End with: 'Learning to prompt AI is like learning to walk. Most people are still crawling — they get generic output and decide AI isn't useful. That's not an AI problem. That's a prompting problem. It's fixable.'\n\nDrop chat pulse check.")
    return slide


def s_data(prs, n):
    slide = blank(prs, DEEP_NAVY)
    chrome(slide, n, "THE DATA")
    headline(slide, "You're only selling for 28% of your week.", fs=28)

    stats = [
        ("28%",  "actual selling time",         "Salesforce 2025",   GOLD),
        ("78%",  "missed quota ↑ from 69%",      "Salesmotion 2025",  RGBColor(220, 80, 80)),
        ("3.7x", "quota with structured AI use", "Gartner 2025",      BLUE),
        ("2x",   "more likely to exceed target", "LinkedIn 2025",     BLUE),
    ]
    cw = Inches(2.9)
    for i, (stat, label, src, color) in enumerate(stats):
        stat_card(slide, Inches(0.45) + i * (cw + Inches(0.17)),
                  Inches(2.0), cw, Inches(3.3),
                  stat, label, src, stat_color=color)

    # Bottom callout
    rr(slide, Inches(0.45), Inches(5.5), Inches(12.4), Inches(0.9), CARD,
       line_color=BLUE_DIM, corner=4000)
    box(slide, Inches(0.65), Inches(5.6), Inches(12.0), Inches(0.7),
        "We're targeting the 3 biggest drains today:  LinkedIn content  ·  Inbox management (21% of your day)  ·  Pre-call research (30–60 min, usually skipped)",
        fn=F_BODY, fs=13, color=OFF_WHITE, align=PP_ALIGN.CENTER)

    notes(slide, "Insert chat beat after 78%: 'Has your company given you formal AI training, or are you figuring this out alone?'\n\nBain 2025: AI could double time spent selling from ~25% to ~50%.")
    return slide


def s_4_layers(prs, n):
    slide = blank(prs)
    chrome(slide, n, "FRAMEWORK")
    headline(slide, "Every prompt has 4 layers. Most people only use 1.", fs=27)

    layers = [
        ("SYSTEM",    GOLD,  "Who the AI is",       "Role, expertise, tone, rules. The job description. Without it: generalist. With it: specialist."),
        ("USER",      BLUE,  "The task",             "Your specific request with full context. More context = better output. Always."),
        ("ASSISTANT", GOLD,  "Quality benchmark",    "Examples of what good looks like. Your best work handed over. This is few-shot prompting — the highest-leverage technique."),
        ("DEVELOPER", DIM,   "Behind the curtain",   "API & automation logic. Used when building agents. Know it exists."),
    ]
    for i, (label, color, sub, desc) in enumerate(layers):
        t = Inches(2.05) + i * Inches(1.25)
        rr(slide, Inches(0.45), t, Inches(12.4), Inches(1.1), CARD,
           line_color=color if color != DIM else CARD, line_w=0.5, corner=4000)
        # Left color strip
        rr(slide, Inches(0.45), t, Inches(0.14), Inches(1.1), color, corner=4000)
        # Pill label
        pill(slide, Inches(0.72), t + Inches(0.32), label, bg=color if color != DIM else DARK_NAVY,
             fg=DEEP_NAVY if color == GOLD else WHITE, fs=10)
        box(slide, Inches(2.0), t + Inches(0.18), Inches(2.5), Inches(0.38),
            sub, fn=F_SUB, fs=13, color=WHITE, bold=True)
        box(slide, Inches(4.7), t + Inches(0.2), Inches(8.0), Inches(0.68),
            desc, fn=F_BODY, fs=13, color=OFF_WHITE)

    rr(slide, Inches(0.45), Inches(7.02), Inches(12.4), Inches(0.32), DARK_NAVY, corner=4000)
    box(slide, Inches(0.45), Inches(7.04), Inches(12.4), Inches(0.28),
        "System = job description  ·  User = daily task  ·  Assistant = quality benchmark",
        fn=F_SUB, fs=12, color=GOLD, align=PP_ALIGN.CENTER)

    notes(slide, "Key: 'Most people type what they want and hit enter. The other three layers are where the results live.'")
    return slide


def s_5_step_loop(prs, n):
    slide = blank(prs)
    chrome(slide, n, "FRAMEWORK")
    headline(slide, "Great prompts are built, not typed.", fs=28)

    steps = [
        ("01", "SET A GOAL",      "Define success before writing anything.", GOLD),
        ("02", "WRITE",           "Get something down fast. Imperfect is fine.", BLUE),
        ("03", "EVALUATE",        "Be specific about what's off.", GOLD),
        ("04", "APPLY TECHNIQUE", "Role. Few-shot. Constraints. One fix at a time.", BLUE),
        ("05", "RE-EVALUATE",     "Did the fix work? Test against the original goal.", GOLD),
    ]
    cw = Inches(2.3)
    for i, (num, title, desc, color) in enumerate(steps):
        l = Inches(0.45) + i * (cw + Inches(0.09))
        rr(slide, l, Inches(2.05), cw, Inches(3.85), CARD,
           line_color=color, line_w=0.5, corner=5000)
        # Number badge
        rr(slide, l + Inches(0.65), Inches(2.25), Inches(1.0), Inches(0.85), color,
           corner=8000)
        box(slide, l + Inches(0.65), Inches(2.25), Inches(1.0), Inches(0.85),
            num, fn=F_HEAD, fs=20, color=DEEP_NAVY, bold=True, align=PP_ALIGN.CENTER)
        box(slide, l + Inches(0.1), Inches(3.3), cw - Inches(0.2), Inches(0.5),
            title, fn=F_SUB, fs=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        box(slide, l + Inches(0.12), Inches(3.88), cw - Inches(0.24), Inches(1.75),
            desc, fn=F_BODY, fs=12, color=OFF_WHITE, align=PP_ALIGN.CENTER)

    rr(slide, Inches(0.45), Inches(6.12), Inches(12.4), Inches(0.42), DARK_NAVY, corner=4000)
    box(slide, Inches(0.45), Inches(6.15), Inches(12.4), Inches(0.35),
        "PromptAnything.io automates steps 3 & 4  —  sign up now, 30 seconds",
        fn=F_SUB, fs=14, color=GOLD, align=PP_ALIGN.CENTER)

    notes(slide, "Key: 'This is a build process, not a template hunt. Stop expecting perfect output on the first try. Start engineering toward it.'")
    return slide


def s_promptanything(prs, n):
    slide = blank(prs)
    chrome(slide, n, "TOOL ORIENTATION")
    headline(slide, "PromptAnything builds the prompt. You approve it.", fs=27)

    steps = [
        ("01", "Drop in rough prompt",      "Plain language. Incomplete is fine.",                               GOLD),
        ("02", "Select: Agent",             "For Claude Projects. Direct use, not automation pipeline.",         BLUE),
        ("03", "Review AI questions",       "Tool surfaces gaps and drafts answers. You review, not write.",     GOLD),
        ("04", "Refine answers",            "Correct what's wrong. Approve what's right.",                      BLUE),
        ("05", "Generate + push to Claude", "Copy → paste as Claude Project instructions → run.",               GOLD),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        t = Inches(2.0) + i * Inches(1.02)
        rr(slide, Inches(0.45), t, Inches(12.4), Inches(0.88), CARD,
           line_color=CARD, corner=4000)
        # Pill number
        rr(slide, Inches(0.6), t + Inches(0.2), Inches(0.5), Inches(0.48), color, corner=8000)
        box(slide, Inches(0.6), t + Inches(0.2), Inches(0.5), Inches(0.48),
            num, fn=F_SUB, fs=11, color=DEEP_NAVY, bold=True, align=PP_ALIGN.CENTER)
        box(slide, Inches(1.3), t + Inches(0.18), Inches(4.5), Inches(0.42),
            title, fn=F_SUB, fs=14, color=WHITE, bold=True)
        box(slide, Inches(6.0), t + Inches(0.2), Inches(6.7), Inches(0.42),
            desc, fn=F_BODY, fs=13, color=OFF_WHITE)
        # subtle rule between steps
        if i < 4:
            shape(slide, RECT, Inches(0.45), t + Inches(0.9), Inches(12.4), Inches(0.01), DARK_NAVY)

    rr(slide, Inches(0.45), Inches(7.1), Inches(12.4), Inches(0.28), DEEP_NAVY, corner=3000)
    box(slide, Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.22),
        "NOW: Create 3 Claude Projects  →  Sales Copywriting Agent  ·  Google Email Hub  ·  Social Chatter Agent",
        fn=F_SUB, fs=12, color=GOLD, align=PP_ALIGN.CENTER)

    notes(slide, "Pause 60-90 sec for sign-ups. Watch chat.\n'Steps 3 and 4 are where this tool earns its value. You are reviewing what it surfaced and approving it. Much faster skill to learn.'")
    return slide


def s_demo_setup(prs, n, section_label_txt, title, pain, rough_prompt, demo_note):
    slide = blank(prs)
    chrome(slide, n, section_label_txt, GOLD)
    headline(slide, title, fs=26)

    # Pain pill
    rr(slide, Inches(0.45), Inches(2.12), Inches(12.4), Inches(0.5), DARK_NAVY,
       line_color=GOLD, line_w=0.5, corner=3000)
    box(slide, Inches(0.65), Inches(2.17), Inches(12.0), Inches(0.38),
        f"⚡  {pain}", fn=F_BODY, fs=13, color=OFF_WHITE)

    # Code editor chrome
    rr(slide, Inches(0.45), Inches(2.78), Inches(12.4), Inches(0.42), RGBColor(40, 45, 65),
       corner=4000)
    shape(slide, RECT, Inches(0.45), Inches(3.1), Inches(12.4), Inches(0.08), RGBColor(40, 45, 65))
    # Traffic light dots
    for j, c in enumerate([RGBColor(255, 90, 90), RGBColor(255, 190, 50), RGBColor(80, 200, 100)]):
        rr(slide, Inches(0.68) + j * Inches(0.28), Inches(2.9), Inches(0.18), Inches(0.18), c, corner=20000)
    box(slide, Inches(1.5), Inches(2.85), Inches(6.0), Inches(0.3),
        "rough_prompt.txt", fn=F_MONO, fs=11, color=DIM)
    pill(slide, Inches(10.5), Inches(2.88), "ROUGH PROMPT", bg=GOLD, fg=DEEP_NAVY, fs=9)

    # Code body
    rr(slide, Inches(0.45), Inches(3.18), Inches(12.4), Inches(3.0),
       RGBColor(12, 16, 32), line_color=RGBColor(40, 50, 75), line_w=0.5, corner=4000)
    shape(slide, RECT, Inches(0.45), Inches(3.18), Inches(12.4), Inches(0.15), RGBColor(12, 16, 32))

    tb = slide.shapes.add_textbox(Inches(0.72), Inches(3.28), Inches(11.9), Inches(2.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text           = rough_prompt
    r.font.name      = F_MONO
    r.font.size      = Pt(13)
    r.font.color.rgb = RGBColor(180, 220, 180)   # green-tinted for code feel

    # Demo cue
    rr(slide, Inches(0.45), Inches(6.35), Inches(12.4), Inches(0.42), DARK_NAVY, corner=3000)
    box(slide, Inches(0.65), Inches(6.4), Inches(12.0), Inches(0.35),
        "Paste into PromptAnything  →  Agent  →  Review questions  →  Refine  →  Generate  →  Push to Claude Project",
        fn=F_SUB, fs=12, color=BLUE, align=PP_ALIGN.CENTER)

    footer_bar(slide)
    notes(slide, demo_note)
    return slide


def s_demo_output(prs, n, section_label_txt, title, points, action, insight=None):
    slide = blank(prs)
    chrome(slide, n, section_label_txt, GOLD)
    headline(slide, title, fs=26)

    for i, pt in enumerate(points):
        t = Inches(2.1) + i * Inches(1.4)
        rr(slide, Inches(0.45), t, Inches(12.4), Inches(1.2), CARD,
           line_color=BLUE_DIM, line_w=0.5, corner=4000)
        # Gold checkmark pill
        rr(slide, Inches(0.65), t + Inches(0.25), Inches(0.65), Inches(0.65), GOLD, corner=8000)
        box(slide, Inches(0.65), t + Inches(0.25), Inches(0.65), Inches(0.65),
            "✓", fn=F_HEAD, fs=16, color=DEEP_NAVY, bold=True, align=PP_ALIGN.CENTER)
        box(slide, Inches(1.5), t + Inches(0.22), Inches(11.2), Inches(0.75),
            pt, fn=F_BODY, fs=15, color=WHITE)

    top_offset = Inches(2.1) + len(points) * Inches(1.4)

    if insight:
        box(slide, Inches(0.45), top_offset + Inches(0.1), Inches(12.4), Inches(0.38),
            f"→  {insight}", fn=F_SUB, fs=13, color=GOLD, italic=True)
        top_offset += Inches(0.52)

    # Action card
    rr(slide, Inches(0.45), top_offset + Inches(0.05), Inches(12.4), Inches(1.3),
       DARK_NAVY, line_color=GOLD, line_w=0.75, corner=4000)
    pill(slide, Inches(0.65), top_offset + Inches(0.18), "YOUR TURN", bg=GOLD, fg=DEEP_NAVY, fs=10)
    box(slide, Inches(0.65), top_offset + Inches(0.56), Inches(12.0), Inches(0.65),
        action, fn=F_BODY, fs=13, color=OFF_WHITE)

    footer_bar(slide)
    return slide


def s_what_built(prs, n):
    slide = blank(prs)
    chrome(slide, n, "WHAT YOU BUILT")
    headline(slide, "Three working systems. Use them tonight.", fs=28)

    systems = [
        (GOLD, "Sales Copywriting Agent",
         "LinkedIn posts in your voice  ·  Prospect comments that open conversations without pitching"),
        (BLUE, "Google Email Hub",
         "Morning briefing from your real inbox before you touch a single email"),
        (GOLD, "Social Chatter Agent",
         "Full research brief on any prospect in under 2 minutes — competitor intel, intent signals, talk track"),
    ]
    for i, (color, name, desc) in enumerate(systems):
        t = Inches(2.05) + i * Inches(1.52)
        rr(slide, Inches(0.45), t, Inches(12.4), Inches(1.28), CARD,
           line_color=color, line_w=0.75, corner=4000)
        shape(slide, RECT, Inches(0.45), t, Inches(0.14), Inches(1.28), color)
        box(slide, Inches(0.78), t + Inches(0.12), Inches(12.0), Inches(0.48),
            name, fn=F_SUB, fs=18, color=color, bold=True)
        box(slide, Inches(0.78), t + Inches(0.65), Inches(11.8), Inches(0.5),
            desc, fn=F_BODY, fs=13, color=OFF_WHITE)

    rr(slide, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.38), DARK_NAVY, corner=3000)
    box(slide, Inches(0.45), Inches(6.78), Inches(12.4), Inches(0.32),
        "Sellers using structured AI: 3.7x more likely to hit quota  ·  Gartner & LinkedIn, 2025",
        fn=F_SUB, fs=13, color=GOLD, align=PP_ALIGN.CENTER)

    notes(slide, "Land each system name. Pause.\n'These are not demos. They're tools you can run tonight — built around YOUR product, YOUR ICP, YOUR voice.'")
    return slide


def s_cta(prs, n):
    slide = blank(prs, DEEP_NAVY)

    # Background depth layer
    shape(slide, RECT, Inches(7.5), Inches(0), Inches(5.83), H, MID_NAVY)
    shape(slide, RECT, Inches(0), Inches(0), Inches(0.22), H, GOLD)

    logo_wordmark(slide)
    footer_bar(slide)
    progress_dots(slide, n)

    headline(slide, "One conversation.", fs=38, t=Inches(0.95))
    headline(slide, "See what's actually possible.", fs=28, color=GOLD, t=Inches(1.72))

    # Divider
    shape(slide, RECT, Inches(0.45), Inches(2.5), Inches(6.5), Inches(0.04), GOLD)

    # Strategy call card
    rr(slide, Inches(0.45), Inches(2.7), Inches(6.5), Inches(1.45), CARD,
       line_color=GOLD, line_w=0.75, corner=5000)
    pill(slide, Inches(0.7), Inches(2.85), "BOOK A CALL", bg=GOLD, fg=DEEP_NAVY, fs=10)
    box(slide, Inches(0.7), Inches(3.22), Inches(6.0), Inches(0.42),
        "calendly.com/richmond", fn=F_SUB, fs=20, color=GOLD, bold=True)
    box(slide, Inches(0.7), Inches(3.66), Inches(6.0), Inches(0.38),
        "30 minutes  ·  No pitch deck  ·  No pressure",
        fn=F_BODY, fs=13, color=OFF_WHITE)

    # Trial card
    rr(slide, Inches(0.45), Inches(4.3), Inches(6.5), Inches(1.3), CARD,
       line_color=BLUE, line_w=0.75, corner=5000)
    pill(slide, Inches(0.7), Inches(4.45), "FREE TRIAL", bg=BLUE, fg=WHITE, fs=10)
    box(slide, Inches(0.7), Inches(4.82), Inches(6.0), Inches(0.42),
        "PromptAnything.io  —  7 days free", fn=F_SUB, fs=18, color=WHITE, bold=True)
    box(slide, Inches(0.7), Inches(5.26), Inches(6.0), Inches(0.35),
        "Everything you built today stays in your account",
        fn=F_BODY, fs=13, color=OFF_WHITE)

    # 24-hour actions
    box(slide, Inches(7.8), Inches(2.55), Inches(5.0), Inches(0.38),
        "NEXT 24 HOURS", fn=F_SUB, fs=12, color=GOLD, bold=True)
    actions = [
        "Customize all 3 prompts with your real product, ICP, voice",
        "Connect one more app in Claude — Calendar, Drive, or Slack",
        "Use the comment generator on a real prospect — actually post it",
    ]
    for i, a in enumerate(actions):
        t = Inches(3.0) + i * Inches(1.1)
        rr(slide, Inches(7.7), t, Inches(5.3), Inches(0.88), CARD,
           line_color=BLUE_DIM, line_w=0.5, corner=4000)
        box(slide, Inches(7.7), t, Inches(0.38), Inches(0.88),
            str(i+1), fn=F_HEAD, fs=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        box(slide, Inches(8.15), t + Inches(0.1), Inches(4.6), Inches(0.65),
            a, fn=F_BODY, fs=12, color=OFF_WHITE)

    # QR placeholder
    rr(slide, Inches(7.7), Inches(6.2), Inches(1.5), Inches(1.5), WHITE,
       line_color=GOLD, line_w=1, corner=3000)
    box(slide, Inches(7.7), Inches(6.95), Inches(1.5), Inches(0.35),
        "QR → Book", fn=F_BODY, fs=9, color=DIM, align=PP_ALIGN.CENTER)

    notes(slide, "CLOSING (90-120 sec):\n'We started with: if AI is supposed to make reps more productive, why are more people missing quota now? You have the answer. It was never about the tools. It was about how you ask.\n\nThe 72% of your week that isn't selling — that's where momentum dies. Three engineered prompts just gave you three ways to get it back.\n\nWhat comes next: autonomous agents, automation pipelines, AI that monitors your pipeline and flags deals before they die quietly. That's what Bishop AI builds.\n\nCalendar link is in the chat. 30 minutes. No pitch deck.\n\nThe trial starts the moment you log in. Use tonight.'\n\n[Hold 3 seconds of silence after final line.]")
    return slide


# ── Build & upload ─────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    s_cover(prs)        # cover

    s_hook(prs, 1)
    s_origin(prs, 2)
    s_data(prs, 3)
    s_4_layers(prs, 4)
    s_5_step_loop(prs, 5)
    s_promptanything(prs, 6)

    s_demo_setup(prs, 7, "BUILD 1 OF 3",
        title="Prompt 1: Sales Copywriting Agent",
        pain="Reps post nothing or sound like everyone else. Prospect comments go unsent or pitch — which is worse than silence.",
        rough_prompt=(
            "I want an AI that writes LinkedIn posts for me in my voice\n"
            "and also writes comments I can post on prospects' LinkedIn posts.\n\n"
            "My name is Amanda. I sell pipeline intelligence software to sales leaders.\n"
            "My tone is direct and insight-led.\n"
            "I don't do inspirational quotes or motivational fluff."
        ),
        demo_note="Refine: Topics (pipeline visibility, deal coaching, AI in sales), Post format (3-5 lines, no hashtag clusters, ends with POV/question), Comment goal (validate + insight + curiosity, never pitch), Voice (peer with firsthand experience).\n\nRun two live prompts:\n1. Post: 78% quota miss, time allocation angle, question for VP of Sales\n2. Comment: Marcus Webb Q1 post — validate, one insight, no pitch")

    s_demo_output(prs, 8, "BUILD 1 OF 3",
        title="Post in your voice. Comment that opens doors.",
        points=[
            "LinkedIn post: real stat + clear POV + question the ICP wants to answer",
            "Prospect comment: validates their post, adds field insight, creates curiosity — no pitch",
        ],
        action="Build yours: your name, product, voice. Find a real prospect post. Run the comment prompt. Drop your output in the chat.",
        insight="That is how a LinkedIn comment becomes a warm first touch.")

    s_demo_setup(prs, 9, "BUILD 2 OF 3",
        title="Prompt 2: Google Email Hub",
        pain="Sales pros spend 21% of the day on email. Most of the inbox is noise hiding the signals that matter. (HubSpot, 2025)",
        rough_prompt=(
            "I want Claude to scan my Gmail inbox every morning and give me a daily briefing.\n\n"
            "Tell me which emails need my attention today,\n"
            "flag anything from clients or active deals,\n"
            "surface calendar issues, and help me filter the noise\n"
            "so I'm not starting my day reactively."
        ),
        demo_note="First: Claude → Settings → Integrations → Gmail → OAuth. Show live.\n\nRefine: Priority senders, urgency definition, calendar flags (conflicts, missing agendas), format (3 sections, 1 page), noise filters.\n\nRun: Morning briefing on live inbox.")

    s_demo_output(prs, 10, "BUILD 2 OF 3",
        title="Your inbox tells you what matters before you open it.",
        points=[
            "Action Required: emails needing reply or task — summary + suggested action",
            "FYI Only: worth knowing, nothing required from you",
            "Calendar Flags: conflicts, missing agendas, meetings that need prep",
        ],
        action="Connect Gmail in Claude. Customize priority senders for your actual clients. Run the briefing. Drop in chat: how many flagged as Action Required?",
        insight="The difference between starting your day on offense vs. defense is usually just this clarity.")

    s_demo_setup(prs, 11, "BUILD 3 OF 3",
        title="Prompt 3: Social Chatter Agent",
        pain="Top reps research every prospect before contact. Manual research: 30-60 min per call. So it gets skipped. (LinkedIn, 2025)",
        rough_prompt=(
            "I want a full research brief on a prospect before I reach out.\n\n"
            "I need: weaknesses of their current tools from real reviews,\n"
            "recent company activity and intent signals,\n"
            "a persona analysis for their specific role,\n"
            "an opening angle for my product,\n"
            "and a talk track opener for a call or LinkedIn message."
        ),
        demo_note="Refine: Product (DealPath — pipeline intelligence), ICP (VP Sales, Dir RevOps at SaaS 100-500), Competitors (Clari, Gong, Salesforce Forecasting — G2/Trustpilot/Reddit), Intent signals, 5-section output.\n\nRun: Full brief on Marcus Webb, VP Sales at Axiom Software.")

    s_demo_output(prs, 12, "BUILD 3 OF 3",
        title="Walk in knowing what their tools are already failing at.",
        points=[
            "Competitor weaknesses from G2, Trustpilot, Reddit — specific, not generic",
            "Intent signals: last 90 days of funding, hiring, leadership changes",
            "Opening angle + talk track opener — ready to send or say right now",
        ],
        action="Build yours: your product, ICP, top 2-3 competitors. Run on a real prospect. Paste the Competitor Weaknesses section in the chat.",
        insight="45 minutes of research → 90 seconds. Every time. For every prospect.")

    s_what_built(prs, 13)
    s_cta(prs, 14)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    return OUTPUT


def upload(path):
    if not UPLOAD.exists():
        print(f"Upload script not found: {UPLOAD}")
        return
    r = subprocess.run(["python", str(UPLOAD), str(path)], capture_output=True, text=True)
    print(r.stdout.strip() if r.returncode == 0 else f"Upload error:\n{r.stderr}")


if __name__ == "__main__":
    print("Building deck...")
    p = build()
    print("Uploading...")
    upload(p)